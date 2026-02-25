"""Extracteur pour fichiers PowerPoint (.pptx)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from app.models import ExtractedContent
from .base import BaseExtractor


class PptxExtractor(BaseExtractor):
    """Extrait le texte des diapositives PowerPoint."""
    
    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]
    
    def extract(self, path: Path) -> ExtractedContent:
        prs = Presentation(str(path))
        sections = []
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                slide_content = "\n".join(slide_texts)
                sections.append({"title": f"Slide {i + 1}", "content": [slide_content]})
                text_parts.append(slide_content)
        
        title = prs.core_properties.title or (f"Présentation ({len(prs.slides)} slides)")
        raw_text = "\n\n---\n\n".join(text_parts)
        return ExtractedContent(raw_text=raw_text, title=title, sections=sections)
